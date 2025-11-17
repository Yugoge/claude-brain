#!/usr/bin/env python3
"""
PowerPoint (PPTX) Parser
Extracts text content from PowerPoint presentations for learning
"""

import sys
import json
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("ERROR: python-pptx not installed")
    print("Install with: pip install python-pptx")
    sys.exit(1)


class PPTParser:
    """Parser for PowerPoint PPTX files."""

    @staticmethod
    def extract_text_from_slide(slide) -> dict:
        """Extract text from a single slide."""
        slide_text = {
            "title": "",
            "content": [],
            "notes": ""
        }

        # Extract title
        if slide.shapes.title:
            slide_text["title"] = slide.shapes.title.text

        # Extract content from text boxes and shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Skip title (already captured)
                if shape == slide.shapes.title:
                    continue
                slide_text["content"].append(shape.text.strip())

        # Extract notes
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide.notes_text_frame:
                slide_text["notes"] = notes_slide.notes_text_frame.text.strip()

        return slide_text

    @staticmethod
    def parse_pptx(file_path: str, output_format: str = "markdown") -> str:
        """
        Parse a PPTX file and extract all text content.

        Args:
            file_path: Path to PPTX file
            output_format: Output format ('markdown', 'json', or 'text')

        Returns:
            Extracted content in specified format
        """
        try:
            prs = Presentation(file_path)
        except Exception as e:
            return f"ERROR: Failed to open {file_path}: {e}"

        # Extract metadata
        metadata = {
            "title": prs.core_properties.title or Path(file_path).stem,
            "author": prs.core_properties.author or "Unknown",
            "total_slides": len(prs.slides),
            "file_path": file_path
        }

        # Extract slides
        slides_data = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_data = PPTParser.extract_text_from_slide(slide)
            slide_data["slide_number"] = idx
            slides_data.append(slide_data)

        # Format output
        if output_format == "json":
            return json.dumps({
                "metadata": metadata,
                "slides": slides_data
            }, indent=2, ensure_ascii=False)

        elif output_format == "markdown":
            return PPTParser._format_as_markdown(metadata, slides_data)

        else:  # plain text
            return PPTParser._format_as_text(metadata, slides_data)

    @staticmethod
    def _format_as_markdown(metadata: dict, slides: list) -> str:
        """Format extracted content as Markdown."""
        lines = []

        # Header
        lines.append(f"# {metadata['title']}\n")
        lines.append(f"**Author**: {metadata['author']}")
        lines.append(f"**Total Slides**: {metadata['total_slides']}")
        lines.append(f"**Source**: `{metadata['file_path']}`\n")
        lines.append("---\n")

        # Slides
        for slide in slides:
            lines.append(f"## Slide {slide['slide_number']}: {slide['title'] or '(No Title)'}\n")

            if slide['content']:
                for content in slide['content']:
                    # Detect bullet points
                    if content.startswith("•") or content.startswith("-"):
                        lines.append(content)
                    else:
                        lines.append(f"{content}\n")

            if slide['notes']:
                lines.append(f"\n**Speaker Notes**:\n> {slide['notes']}\n")

            lines.append("\n---\n")

        return "\n".join(lines)

    @staticmethod
    def _format_as_text(metadata: dict, slides: list) -> str:
        """Format extracted content as plain text."""
        lines = []

        lines.append(f"{metadata['title']}")
        lines.append(f"Author: {metadata['author']}")
        lines.append(f"Total Slides: {metadata['total_slides']}")
        lines.append(f"=" * 60 + "\n")

        for slide in slides:
            lines.append(f"Slide {slide['slide_number']}: {slide['title'] or '(No Title)'}")
            lines.append("-" * 60)

            for content in slide['content']:
                lines.append(content)

            if slide['notes']:
                lines.append(f"\nNotes: {slide['notes']}")

            lines.append("\n")

        return "\n".join(lines)


def main():
    """Command-line interface for PPT parser."""
    if len(sys.argv) < 2:
        print("Usage: python parse-ppt.py <file.pptx> [format]")
        print("Formats: markdown (default), json, text")
        sys.exit(1)

    file_path = sys.argv[1]
    output_format = sys.argv[2] if len(sys.argv) > 2 else "markdown"

    if not Path(file_path).exists():
        print(f"ERROR: File not found: {file_path}")
        sys.exit(1)

    result = PPTParser.parse_pptx(file_path, output_format)
    print(result)


if __name__ == "__main__":
    main()
